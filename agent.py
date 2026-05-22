import io
import json
import re

import streamlit as st
from groq import Groq

# ── File builders ──────────────────────────────────────────────────────────────

NAVY = "1F3864"


def build_excel(data: dict) -> bytes:
    from openpyxl import Workbook
    from openpyxl.styles import Alignment, Font, PatternFill

    wb = Workbook()
    ws = wb.active
    ws.title = data.get("sheet_name", "Sheet1")

    headers = data.get("headers", [])
    rows = data.get("rows", [])

    header_fill = PatternFill(fill_type="solid", fgColor=NAVY)
    header_font = Font(bold=True, color="FFFFFF")

    for col_idx, header in enumerate(headers, start=1):
        cell = ws.cell(row=1, column=col_idx, value=header)
        cell.fill = header_fill
        cell.font = header_font
        cell.alignment = Alignment(horizontal="center")

    for row_idx, row in enumerate(rows, start=2):
        for col_idx, value in enumerate(row, start=1):
            ws.cell(row=row_idx, column=col_idx, value=value)

    # SUM row for numeric columns
    if rows:
        sum_row = len(rows) + 2
        for col_idx, header in enumerate(headers, start=1):
            col_values = []
            for row in rows:
                val = row[col_idx - 1] if col_idx - 1 < len(row) else None
                if isinstance(val, (int, float)):
                    col_values.append(val)
            if col_values:
                col_letter = ws.cell(row=1, column=col_idx).column_letter
                cell = ws.cell(
                    row=sum_row,
                    column=col_idx,
                    value=f"=SUM({col_letter}2:{col_letter}{sum_row - 1})",
                )
                cell.font = Font(bold=True)
            else:
                if col_idx == 1:
                    cell = ws.cell(row=sum_row, column=col_idx, value="Total")
                    cell.font = Font(bold=True)

        for col_idx in range(1, len(headers) + 1):
            ws.column_dimensions[
                ws.cell(row=1, column=col_idx).column_letter
            ].auto_size = True

    buf = io.BytesIO()
    wb.save(buf)
    return buf.getvalue()


def build_powerpoint(data: dict) -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.util import Pt

    prs = Presentation()
    navy_rgb = RGBColor(0x1F, 0x38, 0x64)

    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = navy_rgb

    title_tf = slide.shapes.title.text_frame
    title_tf.text = data.get("title", "Presentation")
    title_tf.paragraphs[0].runs[0].font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    title_tf.paragraphs[0].runs[0].font.size = Pt(40)
    title_tf.paragraphs[0].runs[0].font.bold = True

    subtitle_ph = slide.placeholders[1]
    subtitle_ph.text = data.get("subtitle", "")
    for para in subtitle_ph.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)

    # Content slides
    content_layout = prs.slide_layouts[1]
    for slide_data in data.get("slides", []):
        slide = prs.slides.add_slide(content_layout)

        title_shape = slide.shapes.title
        title_shape.text = slide_data.get("title", "")
        for para in title_shape.text_frame.paragraphs:
            for run in para.runs:
                run.font.color.rgb = navy_rgb
                run.font.bold = True
                run.font.size = Pt(28)

        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        for idx, bullet in enumerate(slide_data.get("bullets", [])):
            para = tf.add_paragraph() if idx > 0 else tf.paragraphs[0]
            para.text = bullet
            para.level = 0

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def build_word(data: dict) -> bytes:
    from docx import Document
    from docx.shared import Pt, RGBColor

    doc = Document()

    title = doc.add_heading(data.get("title", "Document"), level=0)
    title_run = title.runs[0] if title.runs else title.add_run(data.get("title", ""))
    title_run.font.color.rgb = RGBColor(0x1F, 0x38, 0x64)

    for section in data.get("sections", []):
        heading = doc.add_heading(section.get("heading", ""), level=1)
        if heading.runs:
            heading.runs[0].font.color.rgb = RGBColor(0x1F, 0x38, 0x64)
        body = section.get("body", "")
        if body:
            doc.add_paragraph(body)

    buf = io.BytesIO()
    doc.save(buf)
    return buf.getvalue()


# ── LLM prompt ─────────────────────────────────────────────────────────────────

SYSTEM_PROMPT = """You are a file-generation assistant. Given a user request, respond with ONLY valid JSON (no markdown, no explanation).

The JSON must have exactly these fields:
- "file_type": one of "excel", "powerpoint", or "word"
- "filename": a descriptive filename without extension
- "data": an object whose shape depends on file_type

For excel:
{
  "file_type": "excel",
  "filename": "sales_report",
  "data": {
    "sheet_name": "Sales",
    "headers": ["Product", "Q1", "Q2", "Q3", "Q4"],
    "rows": [
      ["Widget A", 1200, 1350, 980, 1500],
      ["Widget B", 870, 920, 1100, 1050]
    ]
  }
}

For powerpoint:
{
  "file_type": "powerpoint",
  "filename": "quarterly_review",
  "data": {
    "title": "Q4 Quarterly Review",
    "subtitle": "FY 2025",
    "slides": [
      {"title": "Key Highlights", "bullets": ["Revenue up 12%", "New markets entered", "Team grew by 20%"]},
      {"title": "Next Steps", "bullets": ["Launch in Europe", "Hire 10 engineers"]}
    ]
  }
}

For word:
{
  "file_type": "word",
  "filename": "project_proposal",
  "data": {
    "title": "Project Proposal",
    "sections": [
      {"heading": "Executive Summary", "body": "This proposal outlines..."},
      {"heading": "Objectives", "body": "Our primary objectives are..."}
    ]
  }
}

Always return only JSON. Never include markdown fences or extra text."""


def call_llm(prompt: str) -> dict:
    client = Groq(api_key=st.secrets["GROQ_API_KEY"])
    response = client.chat.completions.create(
        model="llama-3.1-8b-instant",
        messages=[
            {"role": "system", "content": SYSTEM_PROMPT},
            {"role": "user", "content": prompt},
        ],
        temperature=0.3,
    )
    raw = response.choices[0].message.content.strip()

    # Strip accidental markdown fences
    raw = re.sub(r"^```(?:json)?\s*", "", raw)
    raw = re.sub(r"\s*```$", "", raw)

    return json.loads(raw)


# ── UI ─────────────────────────────────────────────────────────────────────────

EXAMPLES = [
    "Monthly sales report for 5 products across 4 quarters",
    "Investor pitch deck for an AI startup with 5 slides",
    "Business proposal for a mobile app development project",
    "Employee performance review spreadsheet for 8 employees",
    "Marketing strategy document with 4 sections",
]

MIME = {
    "excel": "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "powerpoint": "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "word": "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
}

EXT = {"excel": "xlsx", "powerpoint": "pptx", "word": "docx"}


def main():
    st.set_page_config(page_title="AI File Generator", page_icon="📄", layout="centered")

    st.title("📄 AI File Generator")
    st.caption("Describe what you need — get an Excel, PowerPoint, or Word file instantly.")

    st.markdown("**Quick examples:**")
    cols = st.columns(len(EXAMPLES))
    for col, example in zip(cols, EXAMPLES):
        if col.button(example[:30] + "…", help=example, use_container_width=True):
            st.session_state["prompt"] = example

    st.divider()

    prompt = st.text_area(
        "Describe your file",
        value=st.session_state.get("prompt", ""),
        placeholder="e.g. Quarterly budget spreadsheet for a 5-person startup",
        height=120,
        key="prompt_input",
    )

    generate = st.button("Generate File", type="primary", use_container_width=True)

    if generate:
        if not prompt.strip():
            st.warning("Please enter a prompt before generating.")
            st.stop()

        with st.spinner("Generating your file…"):
            try:
                result = call_llm(prompt)
            except json.JSONDecodeError as e:
                st.error(
                    "The AI returned a response that couldn't be parsed as JSON. "
                    "Try rephrasing your prompt or being more specific."
                )
                with st.expander("Technical details"):
                    st.code(str(e))
                st.stop()
            except Exception as e:
                st.error(f"Something went wrong while calling the AI: {e}")
                st.stop()

        file_type = result.get("file_type", "").lower()
        filename = result.get("filename", "output")
        data = result.get("data", {})

        try:
            if file_type == "excel":
                file_bytes = build_excel(data)
            elif file_type == "powerpoint":
                file_bytes = build_powerpoint(data)
            elif file_type == "word":
                file_bytes = build_word(data)
            else:
                st.error(f"Unknown file type returned by AI: '{file_type}'")
                st.stop()
        except Exception as e:
            st.error(f"Failed to build the {file_type} file: {e}")
            with st.expander("Technical details"):
                st.code(str(e))
            st.stop()

        full_filename = f"{filename}.{EXT[file_type]}"

        st.success(f"Your **{file_type.capitalize()}** file is ready!")
        st.download_button(
            label=f"⬇️ Download {full_filename}",
            data=file_bytes,
            file_name=full_filename,
            mime=MIME[file_type],
            use_container_width=True,
        )

        with st.expander("View raw JSON from AI"):
            st.json(result)


if __name__ == "__main__":
    main()
